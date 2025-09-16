import os
import urllib.parse
import requests
from io import BytesIO
from PIL import Image
from aiogram import Bot, Dispatcher, Router
from aiogram.types import Message, FSInputFile
from aiogram.filters import Command
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from pptx import Presentation
from pptx.util import Inches

# === CONFIG ===
DASHSCOPE_API_KEY = "sk-cb355fe140f84aad88999f6a3db45634"
BOT_TOKEN = "7024550646:AAEse4DfzDfJ7gA4KXIuoGyLpVNAk2-Vp8Q"
API_URL = "https://dashscope-intl.aliyuncs.com/compatible-mode/v1/chat/completions"

# === BOT SETUP ===
storage = MemoryStorage()
bot = Bot(BOT_TOKEN)
dp = Dispatcher(storage=storage)
router = Router()
dp.include_router(router)

# === FSM States ===
class Form(StatesGroup):
    waiting_for_topic = State()
    waiting_for_slide_count = State()


# === START HANDLER ===
@router.message(Command("start"))
async def start_handler(message: Message, state: FSMContext):
    await message.answer("👋 Salom! Menga prezentatsiya mavzusini yuboring:")
    await state.set_state(Form.waiting_for_topic)


# === TOPIC HANDLER ===
@router.message(Form.waiting_for_topic)
async def get_topic(message: Message, state: FSMContext):
    await state.update_data(topic=message.text)
    await message.answer("📄 Nechta slayd kerak? (Masalan: 5)")
    await state.set_state(Form.waiting_for_slide_count)


# === SLIDE COUNT HANDLER ===
@router.message(Form.waiting_for_slide_count)
async def get_slide_count(message: Message, state: FSMContext):
    try:
        count = int(message.text)
    except ValueError:
        await message.answer("❗ Iltimos, faqat raqam kiriting.")
        return

    data = await state.get_data()
    topic = data["topic"]
    await message.answer("⏳ Prezentatsiya tayyorlanmoqda. Iltimos kuting...")

    pptx_file = await create_presentation(topic, count)
    await bot.send_document(message.chat.id, FSInputFile(pptx_file))
    os.remove(pptx_file)

    await state.clear()


# === AI MATN GENERATSIYA ===
def ask_ai(prompt: str) -> str:
    headers = {
        "Authorization": f"Bearer {DASHSCOPE_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "qwen-max",
        "messages": [{"role": "user", "content": prompt}]
    }

    try:
        response = requests.post(API_URL, headers=headers, json=data, timeout=10)
        return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        print("❌ AI xatolik:", e)
        return "AI javob bera olmadi."


# === RASMNI YUKLASH + TEKSHIRISH ===
def download_image_from_unsplash(query: str, save_path: str) -> str:
    try:
        search_url = f"https://source.unsplash.com/800x600/?{urllib.parse.quote(query)}"
        response = requests.get(search_url, allow_redirects=True, timeout=10)

        content_type = response.headers.get("Content-Type", "")
        if not content_type.startswith("image/"):
            print("⚠️ Yuklab olingan fayl rasm emas!")
            return None

        image = Image.open(BytesIO(response.content))
        image.verify()  # Rasmni tekshirish

        with open(save_path, 'wb') as f:
            f.write(response.content)
        return save_path
    except Exception as e:
        print(f"❌ Rasmni yuklashda yoki tekshirishda xatolik: {e}")
        return None


# === PREZENTATSIYA YARATISH ===
async def create_presentation(topic: str, num_slides: int) -> str:
    prs = Presentation()

    for i in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Matn
        prompt = f"{topic} mavzusi bo‘yicha {i+1}-slayd uchun tushunarli va qisqa matn yozing."
        slide_text = ask_ai(prompt)
        title.text = f"{topic} – Slayd {i+1}"
        content.text = slide_text

        # Rasm
        image_path = f"image_{i+1}.jpg"
        if download_image_from_unsplash(topic, image_path) and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(3.5))
                os.remove(image_path)
            except Exception as e:
                print(f"❌ Rasm slaydga qo‘shilmadi: {e}")
        else:
            print(f"❗ {i+1}-slaydga rasm topilmadi, faqat matn qo‘shildi.")

    file_name = f"{topic.replace(' ', '_')}.pptx"
    prs.save(file_name)
    return file_name


# === BOTNI ISHLATISH ===
if __name__ == "__main__":
    import asyncio
    asyncio.run(dp.start_polling(bot))
