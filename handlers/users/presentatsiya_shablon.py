import asyncio
import os
import re
import logging
import tempfile
from pathlib import Path
from io import BytesIO
import textwrap

import aiohttp
from aiogram import Router, F
from aiogram.filters import StateFilter
from aiogram.types import (
    Message,
    FSInputFile,
    BufferedInputFile,
    InlineKeyboardMarkup,
    InlineKeyboardButton,
    CallbackQuery, InputMediaPhoto,
)
from aiogram.fsm.state import StatesGroup, State


from aiogram.fsm.context import FSMContext

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


from loader import bot, db  # <— o‘zingdagi modul nomlari shu bo‘lsa OK

# icrawler for image downloading
from icrawler.builtin import GoogleImageCrawler

logging.basicConfig(level=logging.INFO)

# --- API (sening token va URL qoladi) ---
API_KEY = "sk-cb355fe140f84aad88999f6a3db45634"
API_URL = "https://dashscope-intl.aliyuncs.com/compatible-mode/v1/chat/completions"

router = Router()

# --- Global xotira ---
user_data = {}    # turli sozlamalar (topic, author, bg, image_mode va h.k.)
user_images = {}  # foydalanuvchi yuborgan rasm yo‘llari (temp fayllar)

# -------- STATES --------
class PresState_shablon(StatesGroup):
    language = State()
    topic = State()
    author = State()
    bg = State()
    reja_mode = State()
    reja = State()
    slide_count = State()
    bg_upload = State()
    ozgartirish = State()
    rad_etish=State()
    tanlov=State()
    til=State()
    mavzu=State()
    muallif=State()

class ImageCollectState_shablon(StatesGroup):
    collecting = State()


# =======================
# ====== HELPERS ========
# =======================

def clear_user_images(uid: int, delete_files: bool = True):
    """Foydalanuvchi uchun yig‘ilgan barcha rasm yo‘llarini tozalash."""
    imgs = user_images.get(uid, [])
    if delete_files:
        for p in imgs:
            try:
                if p and os.path.exists(p):
                    os.remove(p)
            except Exception as e:
                logging.warning(f"Temp rasmni o‘chirib bo‘lmadi: {p} | {e}")
    user_images[uid] = []

def download_images(query: str, num_images: int = 1):
    """
    GoogleImageCrawler orqali rasmlarni yuklab oladi va fayl yo'llarini qaytaradi.
    Temp papkaga saqlaydi. (num_images default 1)
    """
    temp_dir = tempfile.mkdtemp(prefix="ppt_images_")
    crawler = GoogleImageCrawler(storage={"root_dir": temp_dir})
    try:
        crawler.crawl(keyword=query, max_num=num_images, min_size=(300, 300))
    except Exception as e:
        logging.warning(f"Image crawl xatosi: {e}")
    image_paths = [
        os.path.join(temp_dir, f) for f in os.listdir(temp_dir)
        if f.lower().endswith(('.jpg', '.jpeg', '.png'))
    ]
    return image_paths

# -------- AI TEXT HELPERS --------
async def ask_ai_content(topic: str, band: str, lan: str, max_retries: int = 4) -> str:
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    if lan == "uz":
        prompt = (
            f"Mavzu: {topic}. Band: {band}. Ushbu band bo'yicha 700-800 belgidan iborat to'liq va izohli matn yozing. "
            "Matn bir nechta to'liq gap va paragraflardan iborat bo'lsin. Matnda -=* kabi belgilar bo'lmasin."
        )
    elif lan == "ru":
        prompt = (
            f"Тема: {topic}. Пункт: {band}. Напишите связанный и пояснительный текст на 700–800 символов по данному пункту. "
            "Текст должен состоять из нескольких предложений и абзацев. Не используйте символы вроде -=*."
        )
    else:
        prompt = (
            f"Topic: {topic}. Section: {band}. Write a complete and explanatory text of 700–800 characters for this section. "
            "The text should consist of several full sentences and paragraphs. Do not use symbols like -=*."
        )

    for attempt in range(max_retries):
        try:
            data = {
                "model": "qwen-turbo",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.7,
                "max_tokens": 1100
            }
            async with aiohttp.ClientSession() as session:
                async with session.post(API_URL, headers=headers, json=data) as resp:
                    result = await resp.json()
                    content = result['choices'][0]['message']['content'].strip()
                    # tozalash
                    content = re.sub(r'[-=*]+', '', content)
                    content = re.sub(r'\n+', '\n', content)
                    # cheklash va tugatish
                    if len(content) < 200:
                        raise ValueError(f"Matn juda qisqa: {len(content)}")
                    if len(content) > 900:
                        content = content[:900]
                        last = max(content.rfind('.'), content.rfind('!'), content.rfind('?'))
                        if last > 0:
                            content = content[:last+1]
                    return content
        except Exception as e:
            logging.warning(f"AI so'rovida xato (urinish {attempt+1}): {e}")
            await asyncio.sleep(1)
    return "Matn olinmadi."

async def ask_ai_content_simple(topic: str, band: str, lan: str) -> str:
    """Fallback – qisqa matn, agar asosiy so‘rov ishlamasa."""
    if lan == "uz":
        return f"{band} mavzusi bo‘yicha qisqacha ma’lumot. {topic} kontekstida tushuntirish va misollar beriladi. Xulosa sifatida asosiy g‘oyalar ko‘rsatiladi."
    if lan == "ru":
        return f"Краткая информация по пункту «{band}». В контексте темы {topic} даются пояснения и примеры. В завершение выделяются ключевые идеи."
    return f"Brief notes for “{band}”. In the context of {topic}, provide explanation and examples. End with the key takeaways."

async def ask_ai_generate_three(topic: str, lan: str) -> list:
    """AI'dan aniq 3 ta asosiy band olish (raqamsiz)."""
    if lan == "uz":
        prompt = f"'{topic}' mavzusi uchun aniq 3 ta asosiy reja bandini yozing. Har bir band yangi qatorda, 3-6 so'zdan iborat bo'lsin. Faqat sarlavha (raqamsiz) qaytaring."
    elif lan == "ru":
        prompt = f"Сгенерируйте ровно 3 основных пункта плана по теме '{topic}'. Каждый пункт с новой строки, 3-6 слов. Возвращайте только заголовки, без нумерации."
    else:
        prompt = f"Generate exactly 3 main outline points for the topic '{topic}'. Each on a new line, 3-6 words. Return only the titles without numbering."

    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    try:
        async with aiohttp.ClientSession() as session:
            data_json = {"model": "qwen-turbo", "messages": [{"role": "user", "content": prompt}], "temperature": 0.6}
            async with session.post(API_URL, headers=headers, json=data_json) as resp:
                result = await resp.json()
                content = result["choices"][0]["message"]["content"]
                lines = [clean_title_line(l) for l in content.split('\n') if l.strip()]
                if len(lines) < 3:
                    lines = [s.strip() for s in re.split(r'\.|\n', content) if s.strip()][:3]
                return lines[:3]
    except Exception as e:
        logging.warning(f"Reja olish xatosi: {e}")
        return [topic, topic, topic]

async def ask_ai_subtopics(topic: str, main_bands: list, needed: int, lan: str) -> list:
    """Kerakli miqdorda subtopic (raqamsiz)."""
    if needed <= 0:
        return []
    mb = "; ".join(main_bands)
    if lan == "uz":
        prompt = (f"'{topic}' mavzusi uchun {needed} ta qisqa (1-6 so'z) alohida subtopic yozing, "
                  f"ular {mb} (asosiy 3ta) ga takrorlanmasin. Har bir yangi qatorda bitta subtopic bo'lsin. Raqamsiz.")
    elif lan == "ru":
        prompt = (f"Для темы '{topic}' с основными пунктами: {mb} сгенерируйте {needed} кратких подтем (1-6 слов), "
                  "которые не повторяют основные пункты. Верните по одной строке на подтему, без нумерации.")
    else:
        prompt = (f"For topic '{topic}' with main points: {mb}, generate {needed} short subtopics (1-6 words) "
                  "that do not duplicate the main points. Return one per line, no numbering.")

    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    try:
        async with aiohttp.ClientSession() as session:
            data_json = {"model": "qwen-turbo", "messages": [{"role": "user", "content": prompt}], "temperature": 0.7}
            async with session.post(API_URL, headers=headers, json=data_json) as resp:
                result = await resp.json()
                content = result["choices"][0]["message"]["content"]
                lines = [clean_title_line(l) for l in content.split('\n') if l.strip()]
                if len(lines) < needed:
                    # Zarurat bo‘lsa, asosiy bandlardan hosila qo‘shamiz
                    extra = []
                    for m in main_bands:
                        extra.append(f"{m} — qo‘shimcha")
                        if len(lines) + len(extra) >= needed:
                            break
                    lines.extend(extra)
                return lines[:needed]
    except Exception as e:
        logging.warning(f"Subtopic olish xatosi: {e}")
        return []

def clean_title_line(line: str) -> str:
    """Raqam va belgilarni yo‘qotish."""
    cleaned = re.sub(r'^\s*\d+[\.\)\-\:]\s*', '', line).strip()
    return cleaned

def format_text(text: str, lan: str, line_length: int = 60) -> str:
    lines = []
    for paragraph in text.split('\n'):
        if paragraph.strip():
            wrapped = textwrap.fill(paragraph, width=line_length)
            lines.append(wrapped)
    return '\n'.join(lines)

# -------- SLIDE HELPERS --------
def add_title_slide(bg_num, prs, bg_bytes, title, author, lan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Agar fon rasmi mavjud bo‘lsa — qo‘shamiz
    if bg_bytes:
        try:
            slide.shapes.add_picture(BytesIO(bg_bytes), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        except Exception as e:
            logging.warning(f"Fon rasm qo‘shishda xato, fonsiz ketadi: {e}")
    white_text = bg_num in (6, 7, 8)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8),
                                         prs.slide_width - Inches(2.4), Inches(1.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = 'Amasis MT Pro Black'
    p.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)

    # Author
    author_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.4),
                                          prs.slide_width - Inches(2.4), Inches(1.0))
    af = author_box.text_frame
    af.clear()
    ap = af.paragraphs[0]
    if lan == 'uz':
        ap.text = f"Muallif: {author}"
    elif lan == 'ru':
        ap.text = f"Автор: {author}"
    else:
        ap.text = f"Author: {author}"
    ap.alignment = PP_ALIGN.CENTER
    ap.font.size = Pt(28)
    ap.font.name = 'Amasis MT Pro Black'
    ap.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)

def add_slide_with_side_image(bg_num, prs, bg_bytes, title, lan, content="", image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_bytes:
        try:
            slide.shapes.add_picture(BytesIO(bg_bytes), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        except Exception as e:
            logging.warning(f"Fon rasm qo‘shishda xato (slide): {e}")
    white_text = bg_num in (6, 7, 8)

    if image_path and os.path.exists(image_path):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5),
                                             prs.slide_width/2 - Inches(1.6), Inches(1))
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7),
                                               prs.slide_width/2 - Inches(1.6),
                                               prs.slide_height - Inches(1.8))
        content_align = PP_ALIGN.LEFT
    else:
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0),
                                             prs.slide_width - Inches(1.6), Inches(1))
        content_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.5),
                                               prs.slide_width - Inches(4.0),
                                               prs.slide_height - Inches(4.0))
        content_align = PP_ALIGN.CENTER

    # Title
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = 'Times New Roman'
    p.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER if not image_path else PP_ALIGN.LEFT

    # Content
    tfc = content_box.text_frame
    tfc.clear()
    if content:
        formatted = format_text(content, lan, line_length=55 if image_path else 80)
        first = True
        for line in formatted.split('\n'):
            if not line.strip():
                continue
            para = tfc.paragraphs[0] if first else tfc.add_paragraph()
            para.text = line
            para.font.size = Pt(16)
            para.font.name = 'Times New Roman'
            para.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)
            para.alignment = content_align
            first = False

    # Rasm
    if image_path and os.path.exists(image_path):
        pic_width = prs.slide_width / 2 - Inches(1.0)
        pic_left = prs.slide_width - pic_width - Inches(0.8)
        pic_top = Inches(1.7)
        try:
            slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width)
        except Exception as e:
            logging.warning(f"Rasm qo‘shishda xato: {e}")


# ==============================
# ===== PRESENTATION CORE ======
# ==============================
async def generate_presentation(uid: int, message: Message, lan):
    """
    - 3 ta asosiy band + kerakli bo‘lsa subtopiclar
    - Rasm rejimi:
        user  -> faqat user yuborgan rasmlar (hech qachon auto qidirmaydi)
        auto  -> Google’dan avtomatik
        none  -> rasmsiz
    - Har 3-slaydda rasm (0,3,6,…).
    """
    data = user_data.get(uid, {})
    topic = data.get("topic", "Mavzu")
    author = data.get("author", "")
    bg_num = int(data.get("bg_num", 1))
    bg = data.get("bg", None)  # bytes
    image_mode = data.get("image_mode", "auto")
    print(image_mode)# 'user' | 'auto' | 'none'

    requested_slide_count = data.get("requested_slide_count", 20)
    try:
        requested_slide_count = int(requested_slide_count)
        if requested_slide_count <= 0:
            requested_slide_count = 20
    except:
        requested_slide_count = 20

    # 1) main 3 band
    main_bands = await ask_ai_generate_three(topic, lan)

    # 2) kerakli bo‘lsa subtopics
    remaining = requested_slide_count - len(main_bands)
    subtopics = []
    if remaining > 0:
        subtopics = await ask_ai_subtopics(topic, main_bands, remaining, lan)

    # Slayd sarlavhalari
    full_titles = []
    for i in range(requested_slide_count):
        if i < len(main_bands):
            full_titles.append(main_bands[i])
        else:
            idx_sub = i - len(main_bands)
            if idx_sub < len(subtopics):
                full_titles.append(subtopics[idx_sub])
            else:
                full_titles.append(main_bands[i % len(main_bands)])

    # Save reja
    user_data.setdefault(uid, {})["reja"] = main_bands
    user_data[uid]["requested_slide_count"] = requested_slide_count

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    add_title_slide(bg_num, prs, bg, topic, author, lan)

    # Reja slide
    reja_text = "\n".join([f"{i+1}. {b}" for i, b in enumerate(main_bands)])
    add_slide_with_side_image(bg_num, prs, bg, ("План" if lan=='ru' else ("Reja" if lan=='uz' else "Plan")), lan, reja_text, image_path=None)

    # Progress
    progress_msg = await message.answer("⏳ Tayyorlanmoqda...")

    # User rasmlar indexi
    if image_mode == "user":
        user_data[uid]["user_img_idx"] = 0
        local_user_imgs = list(user_images.get(uid, []))
    else:
        local_user_imgs = []

    images_cache = {}  # auto rejimda band bo‘yicha cache

    # Slaydlar
    for idx in range(requested_slide_count):
        try:
            slide_title = f"{idx+1}. {clean_title_line(full_titles[idx])}"

            # Qaysi slaydlarga rasm beramiz?
            wants_image = (idx % 3 == 0) and (image_mode != "none")

            image_path = None
            if wants_image:
                if image_mode == "user":
                    # Faqat foydalanuvchi yuborgan rasmlar
                    if local_user_imgs:
                        # Navbat bilan qo‘yamiz
                        img_idx = user_data[uid].get("user_img_idx", 0)
                        if img_idx < len(local_user_imgs):
                            image_path = local_user_imgs[img_idx]
                            user_data[uid]["user_img_idx"] = img_idx + 1
                        else:
                            image_path = None  # rasmlar tugadi -> boshqa rasm qo‘shmaymiz
                    else:
                        image_path = None
                elif image_mode == "auto":
                    band_for_image = f"{topic} {full_titles[idx]}"
                    if band_for_image in images_cache:
                        image_path = images_cache[band_for_image]
                    else:
                        imgs = download_images(band_for_image, num_images=1)
                        image_path = imgs[0] if imgs else None
                        images_cache[band_for_image] = image_path
                else:
                    image_path = None  # none rejim

            # Matn
            try:
                matn = await ask_ai_content(topic, full_titles[idx], lan=lan)
            except Exception:
                matn = await ask_ai_content_simple(topic, full_titles[idx], lan)

            if image_path:
                # qisqartiramiz (chap kolonka)
                short = matn[:420]
                last = max(short.rfind('.'), short.rfind('!'), short.rfind('?'))
                if last > 0:
                    short = short[:last+1]
                matn = short

            add_slide_with_side_image(bg_num, prs, bg, slide_title, lan, matn, image_path=image_path)

            # progress
            try:
                progress = (idx + 1) * 100 // requested_slide_count
                blocks = max(1, min(10, (progress // 10)))
                bar = ''.join(['🟩' if i < blocks else '⬜️' for i in range(10)])
                await progress_msg.edit_text(f"⏳ Tayyorlanmoqda...\n\n{bar} ({progress}%)")
            except Exception:
                pass

        except Exception as e:
            logging.error(f"Xatolik slayd {idx+1}: {e}")
            try:
                fallback_text = await ask_ai_content_simple(topic, full_titles[idx], lan)
            except Exception:
                fallback_text = "Matn mavjud emas."
            add_slide_with_side_image(bg_num, prs, bg, slide_title, lan, fallback_text, image_path=None)
            continue

    # Save and send
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    caption = f"📌 {topic}\n💡 @Slaydai_bot"
    await message.answer_document(
        BufferedInputFile(buffer.read(), filename=f"{topic}.pptx"),
        caption=caption
    )




@router.message(StateFilter(PresState_shablon.topic))
async def get_topic(msg: Message, state: FSMContext):
    olish=db.select_shablon_pre(tg_id=msg.from_user.id,til='uz')
    user_data[msg.from_user.id] = {"topic": msg.text}
    user_data[msg.from_user.id]["author"] = olish[3]
    user_data.setdefault(msg.from_user.id, {})["bg_num"] = int(olish[2])
    await state.update_data({'lang': olish[6]})
    rasim_manzili = Path(rf".\handlers\users\rasimlar\{olish[2]}.jpg")
    with open(rasim_manzili, "rb") as f:
        user_data[msg.from_user.id]["bg"] = f.read()
    await state.update_data(reja_mode="ai")
    user_data.setdefault(msg.from_user.id, {})["requested_slide_count"] = olish[5]
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✅ Rasm tanlash", callback_data="img_select_shab")],
        [InlineKeyboardButton(text="♻️ Avtomatik rasmlar", callback_data="img_auto_shab")],
        [InlineKeyboardButton(text="🚫 Rasmsiz", callback_data="img_none_shab")]
    ])

    a=await msg.answer(
        "✅ *Rasm tanlash* — mavzuga oid rasmlarni o‘zingiz yuborasiz (7 tagacha) va faqat shu rasmlar ishlatiladi.\n"
        "♻️ *Avtomatik rasmlar* — bot mavzuga oid rasmlarni internetdan yuklaydi.\n"
        "🚫 *Rasmsiz* — taqdimot rasm ishlatmaydi.",
        reply_markup=kb,
        parse_mode="Markdown"
    )
    await state.update_data(m_id=a.message_id)
    await state.update_data(reja_mode="ai")


@router.callback_query(F.data == "img_select_shab")
async def select_images(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol.get("m_id")
    await bot.delete_message(chat_id=callback.message.chat.id, message_id=m_id)
    uid = callback.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "user"
    clear_user_images(uid, delete_files=True)

    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✅ Tugatish", callback_data="img_finish")],
        [InlineKeyboardButton(text="♻ Rasmlarni 0lash", callback_data="img_reset")]
    ])
    await callback.message.answer(
        "📤 Mavzuga oid rasmlarni yuboring (7 tagacha).\n"
        "Tayyor bo‘lsangiz — *Tugatish* tugmasini bosing.\n"
        "Adashsangiz — *Rasmlarni 0lash* bilan qayta boshlang.\n"
        "🌃 Rasmlar: 0/7",
        reply_markup=kb,
        parse_mode="Markdown"
    )
    await state.set_state(ImageCollectState_shablon.collecting)

@router.message(StateFilter(ImageCollectState_shablon.collecting), F.photo)
async def collect_images(msg: Message, state: FSMContext):
    uid = msg.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "user"
    if uid not in user_images:
        user_images[uid] = []
    if len(user_images[uid]) >= 7:
        await msg.answer("❌ Siz allaqachon 7 ta rasm yubordingiz!")
        return

    file = await bot.get_file(msg.photo[-1].file_id)
    img = await bot.download_file(file.file_path)
    tmp_path = tempfile.mktemp(suffix=".jpg")
    with open(tmp_path, "wb") as f:
        f.write(img.read())

    user_images[uid].append(tmp_path)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✅ Tugatish", callback_data="img_finish_shab")],
        [InlineKeyboardButton(text="♻ Rasmlarni 0lash", callback_data="img_reset_shab")]
    ])
    await msg.answer(f"🌃 Rasmlar: {len(user_images[uid])}/7", reply_markup=kb)

@router.callback_query(F.data == "img_reset_shab")
async def reset_images(callback: CallbackQuery):

    uid = callback.from_user.id
    # barcha yuborgan rasmlarini bekor qilamiz (fayllarni ham o‘chiramiz)
    clear_user_images(uid, delete_files=True)
    await callback.message.answer(
        "♻ Rasmlar 0 qilindi. Qaytadan yuborishingiz mumkin.\n🌃 Rasmlar: 0/7"
    )

@router.callback_query(F.data == "img_finish_shab")
async def finish_images(callback: CallbackQuery, state: FSMContext):
    uid = callback.from_user.id
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="✅ Tayyorlash", callback_data="tayyorlash_shab"),
                InlineKeyboardButton(text="🚫 Rad etish", callback_data="rad_etish_shab"),
                InlineKeyboardButton(text="✏️ O'zgartirish", callback_data="ozgartirish_shab")
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a = await bot.send_photo(chat_id=callback.from_user.id, caption=
    "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
    f"Mavzu: {topic}\n"
    f"Muallif: {author}\n"
    f"Sahifalar soni: {requested_slide_count}\n"
    f"Tili: {lan}\n\n"
    f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
    f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
    f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing",
                             reply_markup=kb, photo=rasim
                             )
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState_shablon.tanlov)
    await callback.answer()

# --- Rasmsiz rejim ---
@router.callback_query(F.data == "img_none_shab")
async def no_images(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol.get("m_id")
    await bot.delete_message(chat_id=callback.message.chat.id, message_id=m_id)
    uid = callback.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "none"
    clear_user_images(uid, delete_files=True)  # albatta ishlatilmaydi
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="✅ Tayyorlash", callback_data="tayyorlash_shab"),
                InlineKeyboardButton(text="🚫 Rad etish", callback_data="rad_etish_shab"),
                InlineKeyboardButton(text="✏️ O'zgartirish", callback_data="ozgartirish_shab")
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a = await bot.send_photo(chat_id=callback.from_user.id, caption=
    "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
    f"Mavzu: {topic}\n"
    f"Muallif: {author}\n"
    f"Sahifalar soni: {requested_slide_count}\n"
    f"Tili: {lan}\n\n"
    f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
    f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
    f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing",
                             reply_markup=kb, photo=rasim
                             )
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState_shablon.tanlov)
    await callback.answer()





# --- Avtomatik rasm rejimi ---
@router.callback_query(F.data == "img_auto_shab")
async def img_auto_handler(callback: CallbackQuery, state: FSMContext):
    uid = callback.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "auto"
    clear_user_images(uid, delete_files=True)
    ol=await state.get_data()
    m_id=ol.get("m_id")
    await bot.delete_message(chat_id=callback.message.chat.id, message_id=m_id)
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")


    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="✅ Tayyorlash", callback_data="tayyorlash_shab"),
                InlineKeyboardButton(text="🚫 Rad etish", callback_data="rad_etish_shab"),
                InlineKeyboardButton(text="✏️ O'zgartirish", callback_data="ozgartirish_shab")
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a=await bot.send_photo(chat_id=callback.from_user.id,caption=
        "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
        f"Mavzu: {topic}\n"
        f"Muallif: {author}\n"
        f"Sahifalar soni: {requested_slide_count}\n"
        f"Tili: {lan}\n\n"
        f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
        f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
        f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing",
        reply_markup=kb,photo=rasim
    )
    await state.update_data({"m_id":a.message_id})
    await state.set_state(PresState_shablon.tanlov)
    await callback.answer()



@router.callback_query(F.data == "ozgartirish_shab", PresState_shablon.tanlov)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol.get("m_id")
    uid=callback.from_user.id
    await bot.delete_message(chat_id=callback.from_user.id, message_id=m_id)
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")


    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu_shab"),
                InlineKeyboardButton(text="Muallif ", callback_data="muallif_shab"),
                InlineKeyboardButton(text="Til", callback_data="til_shab")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus_shab"),
                InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}", callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus_shab")
            ],
            [
                InlineKeyboardButton(text="⬅️", callback_data="oldingi_shab"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="➡️", callback_data="keyingi_shab")
            ],
            [
                InlineKeyboardButton(text="Saqlash ✅", callback_data="saqlash_shab"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a=await bot.send_photo(chat_id=callback.from_user.id,caption=
        "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
        f"Mavzu: {topic}\n"
        f"Muallif: {author}\n"
        f"Sahifalar soni: {requested_slide_count}\n"
        f"Tili: {lan}\n\n"
        f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
        f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
        f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing",
        reply_markup=kb,photo=rasim
    )
    await state.update_data({"m_id":a.message_id})
    await state.set_state(PresState_shablon.ozgartirish)



@router.callback_query(F.data=="mavzu_shab",PresState_shablon.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    a = await callback.message.answer("Taqdimot mavzusini yuboring:")
    await state.update_data({"mm_id": a.message_id})

    await state.set_state(PresState_shablon.mavzu)
    await callback.answer()
    print(2)
@router.callback_query(F.data=="til_shab",PresState_shablon.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🇺🇿 O'zbekcha", callback_data="uz")],
        [InlineKeyboardButton(text="🇷🇺 Русский", callback_data="ru")],
        [InlineKeyboardButton(text="🇬🇧 English", callback_data="en")]
    ])
    a = await callback.message.answer("Taqdimot tilini tanlang:", reply_markup=keyboard)
    await state.update_data({"mm_id": a.message_id})
    await state.set_state(PresState_shablon.til)
    await callback.answer()
    print(1)


@router.callback_query(F.data=="muallif_shab",PresState_shablon.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):

        a = await callback.message.answer("Muallif ismini yuboring:")
        await state.update_data({"mm_id": a.message_id})
        await state.set_state(PresState_shablon.muallif)
        await callback.answer()
        print(3)


@router.callback_query(PresState_shablon.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    """
    Foydalanuvchi shablon parametrlarini o'zgartirish uchun handler.
    Barcha tugmalar (+, -, saqlash, oldingi/keyingi) uchun ishlaydi.
    """
    data = callback.data
    uid = callback.from_user.id

    if data == 'saqlash_shab':
        # SAQLASH tugmasi bosilganda
        ol = await state.get_data()
        m_id = ol.get("m_id")

        # Eski xabarni o'chirish
        await bot.delete_message(chat_id=uid, message_id=m_id)

        # Ma'lumotlarni olish
        user_lang = user_data[uid].get("lang", 'uz')
        author = user_data[uid].get("author")
        requested_slide_count = user_data[uid].get("requested_slide_count", 0)
        bg_num = user_data[uid].get("bg_num")

        # Prezentatsiya generatsiya qilish
        await generate_presentation(uid, callback.message, user_lang)

        # Shablonni bazaga saqlash (hozir commentda)
        # db.add_shablon_pre(
        #     tg_id=uid,
        #     bg_num=bg_num,
        #     ism_fam=author,
        #     slayid_soni=requested_slide_count,
        #     til='uz',
        #     pre_tili=user_lang
        # )

        await state.clear()

    else:
        # Boshqa tugmalar bosilganda
        ol = await state.get_data()
        m_id = ol.get("m_id")
        bg_num = user_data[uid].get("bg_num")
        requested_slide_count = user_data[uid].get("requested_slide_count", 0)

        # + TUGMASI (Slaydlar sonini oshirish)
        if data == "plus_shab":
            if int(requested_slide_count) > 19:
                await callback.answer("Eng ko'pi bilan 20 ta sahifa tanlashingiz mumkin.❗️", show_alert=True)
                await state.set_state(PresState_shablon.ozgartirish)
            else:
                user_data.setdefault(uid, {})["requested_slide_count"] = int(requested_slide_count) + 1
                await state.set_state(PresState_shablon.ozgartirish)

        # - TUGMASI (Slaydlar sonini kamaytirish)
        elif data == "minus_shab":
            if int(requested_slide_count) < 6:
                await callback.answer("Eng kamida 5 ta sahifa tanlashingiz mumkin.❗️", show_alert=True)
                await state.set_state(PresState_shablon.ozgartirish)
            else:
                user_data.setdefault(uid, {})["requested_slide_count"] = int(requested_slide_count) - 1
                await state.set_state(PresState_shablon.ozgartirish)

        # OLINGI TUGMASI (Dizaynni o'zgartirish)
        elif data == "oldingi_shab":
            if int(bg_num) == 1:
                new_bg_num = 10
            else:
                new_bg_num = int(bg_num) - 1

            user_data.setdefault(uid, {})["bg_num"] = new_bg_num
            image_path = Path(rf".\handlers\users\rasimlar\{new_bg_num}.jpg")
            await state.set_state(PresState_shablon.ozgartirish)

            with open(image_path, "rb") as f:
                user_data[uid]["bg"] = f.read()

        # KEYINGI TUGMASI (Dizaynni o'zgartirish)
        elif data == "keyingi_shab":
            if int(bg_num) == 10:
                new_bg_num = 1
            else:
                new_bg_num = int(bg_num) + 1

            user_data.setdefault(uid, {})["bg_num"] = new_bg_num
            image_path = Path(rf".\handlers\users\rasimlar\{new_bg_num}.jpg")
            await state.set_state(PresState_shablon.ozgartirish)

            with open(image_path, "rb") as f:
                user_data[uid]["bg"] = f.read()

        # Yangilangan ma'lumotlarni olish
        state_data = await state.get_data()
        user_lang = state_data.get('lang', 'uz')
        topic = user_data[uid].get("topic")
        author = user_data[uid].get("author")
        requested_slide_count = user_data[uid].get("requested_slide_count", 0)
        bg_num = user_data[uid].get("bg_num")

        # Klaviatura yaratish
        kb = InlineKeyboardMarkup(
            inline_keyboard=[
                [
                    InlineKeyboardButton(text="Mavzu", callback_data="mavzu_shab"),
                    InlineKeyboardButton(text="Muallif", callback_data="muallif_shab"),
                    InlineKeyboardButton(text="Til", callback_data="til_shab")
                ],
                [
                    InlineKeyboardButton(text="-", callback_data="minus_shab"),
                    InlineKeyboardButton(
                        text=f"Sahifalar: {requested_slide_count}",
                        callback_data=f"count:{requested_slide_count}"
                    ),
                    InlineKeyboardButton(text="+", callback_data="plus_shab")
                ],
                [
                    InlineKeyboardButton(text="⬅️", callback_data="oldingi_shab"),
                    InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                    InlineKeyboardButton(text="➡️", callback_data="keyingi_shab")
                ],
                [
                    InlineKeyboardButton(text="Saqlash ✅", callback_data="saqlash_shab"),
                ]
            ]
        )

        # Yangi rasm yuborish
        photo_path = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
        rasim = FSInputFile(photo_path)

        media = InputMediaPhoto(
            media=rasim,
            caption=(
                "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
                f"Mavzu: {topic}\n"
                f"Muallif: {author}\n"
                f"Sahifalar soni: {requested_slide_count}\n"
                f"Tili: {user_lang}\n\n"
                "• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
                "• O'zgartirish uchun yuqoridagi tugmalardan foydalaning"
            )
        )

        # Xabarni yangilash
        updated_msg = await bot.edit_message_media(
            chat_id=uid,
            media=media,
            message_id=m_id,
            reply_markup=kb
        )

        await state.update_data({"m_id": updated_msg.message_id})
        await state.set_state(PresState_shablon.ozgartirish)





@router.callback_query(F.data == "rad_etish_shab", PresState_shablon.tanlov)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    await bot.send_message(chat_id=callback.from_user.id,text="Taqdimot bekor qilindi.\n"
                                "Yangi taqdimot yaratish uchun quyidagi 📕Taqdimot tugmasini bosing!\n\n"
                                "Boshqa amallarni bajarish uchun /start tugmasini bosing.")

    await state.clear()
    await callback.answer()


@router.callback_query(StateFilter(PresState_shablon.til))
async def handle_mavzu_message(message: CallbackQuery, state: FSMContext):
    til=message.data
    print(til)
    ol = await state.get_data()
    mm_id = ol.get("mm_id")
    await bot.delete_message(chat_id=message.from_user.id, message_id=mm_id)
    db.update_shablon_lang(tg_id=message.from_user.id,pre_tili='uz' ,til=message.data)
    user_data.setdefault(message.from_user.id, {})["lan"]=message.data
    ol = await state.get_data()
    m_id = ol.get("m_id")
    uid = message.from_user.id

    data = await state.get_data()
    lan = message.data
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")

    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu_shab"),
                InlineKeyboardButton(text="Muallif ", callback_data="muallif_shab"),
                InlineKeyboardButton(text="Til", callback_data="til_shab")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus_shab"),
                InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}",
                                     callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus_shab")
            ],
            [
                InlineKeyboardButton(text="⬅️", callback_data="oldingi_shab"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="➡️", callback_data="keyingi_shab")
            ],
            [
                InlineKeyboardButton(text="Saqlash ✅", callback_data="saqlash_shab"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")

    rasim = FSInputFile(photo_manzili)
    media = InputMediaPhoto(
        media=rasim, caption=
        "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
        f"Mavzu: {topic}\n"
        f"Muallif: {author}\n"
        f"Sahifalar soni: {requested_slide_count}\n"
        f"Tili: {lan}\n\n"
        f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
        f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
        f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing")
    # await bot.delete_message(chat_id=message.from_user.id, message_id=m_id)
    # a=await bot.send_photo(chat_id=message.from_user.id,photo=rasim,caption=
    #     "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
    #     f"Mavzu: {topic}\n"
    #     f"Muallif: {author}\n"
    #     f"Sahifalar soni: {requested_slide_count}\n"
    #     f"Tili: {lan}\n\n"
    #     f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
    #     f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
    #     f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing",reply_markup=kb)
    a = await bot.edit_message_media(chat_id=message.from_user.id, media=media, message_id=m_id, reply_markup=kb)
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState_shablon.ozgartirish)





@router.message(StateFilter(PresState_shablon.muallif))
async def handle_mavzu_message(message: Message, state: FSMContext):

    ol = await state.get_data()
    mm_id=ol.get("mm_id")
    await bot.delete_message(chat_id=message.from_user.id,message_id=mm_id)
    await message.delete()
    user_data.setdefault(message.from_user.id, {})["author"] = message.text
    ol = await state.get_data()
    m_id = ol.get("m_id")
    uid = message.from_user.id

    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu_shab"),
                InlineKeyboardButton(text="Muallif ", callback_data="muallif_shab"),
                InlineKeyboardButton(text="Til", callback_data="til_shab")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus_shab"),
                InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}",
                                     callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus_shab")
            ],
            [
                InlineKeyboardButton(text="⬅️", callback_data="oldingi_shab"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="➡️", callback_data="keyingi_shab")
            ],
            [
                InlineKeyboardButton(text="Saqlash ✅", callback_data="saqlash_shab"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")

    rasim = FSInputFile(photo_manzili)
    media = InputMediaPhoto(
        media=rasim, caption=
        "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
        f"Mavzu: {topic}\n"
        f"Muallif: {author}\n"
        f"Sahifalar soni: {requested_slide_count}\n"
        f"Tili: {lan}\n\n"
        f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
        f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
        f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing")

    a = await bot.edit_message_media(chat_id=message.from_user.id, media=media, message_id=m_id, reply_markup=kb)
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState_shablon.ozgartirish)










@router.message(StateFilter(PresState_shablon.mavzu))
async def handle_mavzu_message(message: Message, state: FSMContext):

    ol = await state.get_data()
    mm_id=ol.get("mm_id")
    await bot.delete_message(chat_id=message.from_user.id,message_id=mm_id)
    await message.delete()
    user_data.setdefault(message.from_user.id, {})["topic"] = message.text
    ol = await state.get_data()
    m_id = ol.get("m_id")
    uid = message.from_user.id

    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu_shab"),
                InlineKeyboardButton(text="Muallif ", callback_data="muallif_shab"),
                InlineKeyboardButton(text="Til", callback_data="til_shab")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus_shab"),
                InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}",
                                     callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus_shab")
            ],
            [
                InlineKeyboardButton(text="⬅️", callback_data="oldingi_shab"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="➡️", callback_data="keyingi_shab")
            ],
            [
                InlineKeyboardButton(text="Saqlash ✅", callback_data="saqlash_shab"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")

    rasim = FSInputFile(photo_manzili)
    media = InputMediaPhoto(
        media=rasim, caption=
        "🌟Ajoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
        f"Mavzu: {topic}\n"
        f"Muallif: {author}\n"
        f"Sahifalar soni: {requested_slide_count}\n"
        f"Tili: {lan}\n\n"
        f"• Ma'lumotlar to'g'ri bo'lsa, '✅ Tayyorlash' tugmasini bosing!\n"
        f"• O'zgartirish uchun '✏️ O'zgartirish' tugmasini bosing\n"
        f"• Bekor qilish uchun '🚫 Rad etish' tugmasini bosing")

    a = await bot.edit_message_media(chat_id=message.from_user.id, media=media, message_id=m_id, reply_markup=kb)
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState_shablon.ozgartirish)






@router.callback_query(F.data == "tayyorlash_shab", PresState_shablon.tanlov)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol["m_id"]
    await bot.delete_message(chat_id=callback.from_user.id,message_id=m_id)
    uid = callback.from_user.id

    data = await state.get_data()
    lan = user_data[uid].get("lang", 'uz')

    author = user_data[uid].get("author")

    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")
    await generate_presentation(uid, callback.message, lan)



    await state.clear()



