import asyncio
import os
import re
import logging
import tempfile
from pathlib import Path
from io import BytesIO
import textwrap

import aiohttp
from aiogram import Router, F
from aiogram.filters import StateFilter
from aiogram.types import (
    Message,
    FSInputFile,
    BufferedInputFile,
    InlineKeyboardMarkup,
    InlineKeyboardButton,
    CallbackQuery, InputMediaPhoto,
)
from aiogram.fsm.state import StatesGroup, State
from aiogram.fsm.context import FSMContext

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from handlers.users.presentatsiya_shablon import PresState_shablon
from loader import bot, db  # <‚Äî o‚Äòzingdagi modul nomlari shu bo‚Äòlsa OK

# icrawler for image downloading
from icrawler.builtin import GoogleImageCrawler

logging.basicConfig(level=logging.INFO)

# --- API (sening token va URL qoladi) ---
API_KEY = "sk-cb355fe140f84aad88999f6a3db45634"
API_URL = "https://dashscope-intl.aliyuncs.com/compatible-mode/v1/chat/completions"

router = Router()

# --- Global xotira ---
user_data = {}    # turli sozlamalar (topic, author, bg, image_mode va h.k.)
user_images = {}  # foydalanuvchi yuborgan rasm yo‚Äòllari (temp fayllar)

# -------- STATES --------
class PresState(StatesGroup):
    language = State()
    topic = State()
    author = State()
    bg = State()
    reja_mode = State()
    reja = State()
    slide_count = State()
    bg_upload = State()
    tayyorlash=State()
    ozgartirish=State()
    rad_etish=State()
    tanlov=State()
    til=State()
    mavzu=State()
    muallif=State()

class ImageCollectState(StatesGroup):
    collecting = State()


# =======================
# ====== HELPERS ========
# =======================

def clear_user_images(uid: int, delete_files: bool = True):
    """Foydalanuvchi uchun yig‚Äòilgan barcha rasm yo‚Äòllarini tozalash."""
    imgs = user_images.get(uid, [])
    if delete_files:
        for p in imgs:
            try:
                if p and os.path.exists(p):
                    os.remove(p)
            except Exception as e:
                logging.warning(f"Temp rasmni o‚Äòchirib bo‚Äòlmadi: {p} | {e}")
    user_images[uid] = []

def download_images(query: str, num_images: int = 1):
    """
    GoogleImageCrawler orqali rasmlarni yuklab oladi va fayl yo'llarini qaytaradi.
    Temp papkaga saqlaydi. (num_images default 1)
    """
    temp_dir = tempfile.mkdtemp(prefix="ppt_images_")
    crawler = GoogleImageCrawler(storage={"root_dir": temp_dir})
    try:
        crawler.crawl(keyword=query, max_num=num_images, min_size=(300, 300))
    except Exception as e:
        logging.warning(f"Image crawl xatosi: {e}")
    image_paths = [
        os.path.join(temp_dir, f) for f in os.listdir(temp_dir)
        if f.lower().endswith(('.jpg', '.jpeg', '.png'))
    ]
    return image_paths

# -------- AI TEXT HELPERS --------
async def ask_ai_content(topic: str, band: str, lan: str, max_retries: int = 4) -> str:
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    if lan == "uz":
        prompt = (
            f"Mavzu: {topic}. Band: {band}. Ushbu band bo'yicha 700-800 belgidan iborat to'liq va izohli matn yozing. "
            "Matn bir nechta to'liq gap va paragraflardan iborat bo'lsin. Matnda -=* kabi belgilar bo'lmasin."
        )
    elif lan == "ru":
        prompt = (
            f"–¢–µ–º–∞: {topic}. –ü—É–Ω–∫—Ç: {band}. –ù–∞–ø–∏—à–∏—Ç–µ —Å–≤—è–∑–∞–Ω–Ω—ã–π –∏ –ø–æ—è—Å–Ω–∏—Ç–µ–ª—å–Ω—ã–π —Ç–µ–∫—Å—Ç –Ω–∞ 700‚Äì800 —Å–∏–º–≤–æ–ª–æ–≤ –ø–æ –¥–∞–Ω–Ω–æ–º—É –ø—É–Ω–∫—Ç—É. "
            "–¢–µ–∫—Å—Ç –¥–æ–ª–∂–µ–Ω —Å–æ—Å—Ç–æ—è—Ç—å –∏–∑ –Ω–µ—Å–∫–æ–ª—å–∫–∏—Ö –ø—Ä–µ–¥–ª–æ–∂–µ–Ω–∏–π –∏ –∞–±–∑–∞—Ü–µ–≤. –ù–µ –∏—Å–ø–æ–ª—å–∑—É–π—Ç–µ —Å–∏–º–≤–æ–ª—ã –≤—Ä–æ–¥–µ -=*."
        )
    else:
        prompt = (
            f"Topic: {topic}. Section: {band}. Write a complete and explanatory text of 700‚Äì800 characters for this section. "
            "The text should consist of several full sentences and paragraphs. Do not use symbols like -=*."
        )

    for attempt in range(max_retries):
        try:
            data = {
                "model": "qwen-turbo",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.7,
                "max_tokens": 1100
            }
            async with aiohttp.ClientSession() as session:
                async with session.post(API_URL, headers=headers, json=data) as resp:
                    result = await resp.json()
                    content = result['choices'][0]['message']['content'].strip()
                    # tozalash
                    content = re.sub(r'[-=*]+', '', content)
                    content = re.sub(r'\n+', '\n', content)
                    # cheklash va tugatish
                    if len(content) < 200:
                        raise ValueError(f"Matn juda qisqa: {len(content)}")
                    if len(content) > 900:
                        content = content[:900]
                        last = max(content.rfind('.'), content.rfind('!'), content.rfind('?'))
                        if last > 0:
                            content = content[:last+1]
                    return content
        except Exception as e:
            logging.warning(f"AI so'rovida xato (urinish {attempt+1}): {e}")
            await asyncio.sleep(1)
    return "Matn olinmadi."

async def ask_ai_content_simple(topic: str, band: str, lan: str) -> str:
    """Fallback ‚Äì qisqa matn, agar asosiy so‚Äòrov ishlamasa."""
    if lan == "uz":
        return f"{band} mavzusi bo‚Äòyicha qisqacha ma‚Äôlumot. {topic} kontekstida tushuntirish va misollar beriladi. Xulosa sifatida asosiy g‚Äòoyalar ko‚Äòrsatiladi."
    if lan == "ru":
        return f"–ö—Ä–∞—Ç–∫–∞—è –∏–Ω—Ñ–æ—Ä–º–∞—Ü–∏—è –ø–æ –ø—É–Ω–∫—Ç—É ¬´{band}¬ª. –í –∫–æ–Ω—Ç–µ–∫—Å—Ç–µ —Ç–µ–º—ã {topic} –¥–∞—é—Ç—Å—è –ø–æ—è—Å–Ω–µ–Ω–∏—è –∏ –ø—Ä–∏–º–µ—Ä—ã. –í –∑–∞–≤–µ—Ä—à–µ–Ω–∏–µ –≤—ã–¥–µ–ª—è—é—Ç—Å—è –∫–ª—é—á–µ–≤—ã–µ –∏–¥–µ–∏."
    return f"Brief notes for ‚Äú{band}‚Äù. In the context of {topic}, provide explanation and examples. End with the key takeaways."

async def ask_ai_generate_three(topic: str, lan: str) -> list:
    """AI'dan aniq 3 ta asosiy band olish (raqamsiz)."""
    if lan == "uz":
        prompt = f"'{topic}' mavzusi uchun aniq 3 ta asosiy reja bandini yozing. Har bir band yangi qatorda, 3-6 so'zdan iborat bo'lsin. Faqat sarlavha (raqamsiz) qaytaring."
    elif lan == "ru":
        prompt = f"–°–≥–µ–Ω–µ—Ä–∏—Ä—É–π—Ç–µ —Ä–æ–≤–Ω–æ 3 –æ—Å–Ω–æ–≤–Ω—ã—Ö –ø—É–Ω–∫—Ç–∞ –ø–ª–∞–Ω–∞ –ø–æ —Ç–µ–º–µ '{topic}'. –ö–∞–∂–¥—ã–π –ø—É–Ω–∫—Ç —Å –Ω–æ–≤–æ–π —Å—Ç—Ä–æ–∫–∏, 3-6 —Å–ª–æ–≤. –í–æ–∑–≤—Ä–∞—â–∞–π—Ç–µ —Ç–æ–ª—å–∫–æ –∑–∞–≥–æ–ª–æ–≤–∫–∏, –±–µ–∑ –Ω—É–º–µ—Ä–∞—Ü–∏–∏."
    else:
        prompt = f"Generate exactly 3 main outline points for the topic '{topic}'. Each on a new line, 3-6 words. Return only the titles without numbering."

    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    try:
        async with aiohttp.ClientSession() as session:
            data_json = {"model": "qwen-turbo", "messages": [{"role": "user", "content": prompt}], "temperature": 0.6}
            async with session.post(API_URL, headers=headers, json=data_json) as resp:
                result = await resp.json()
                content = result["choices"][0]["message"]["content"]
                lines = [clean_title_line(l) for l in content.split('\n') if l.strip()]
                if len(lines) < 3:
                    lines = [s.strip() for s in re.split(r'\.|\n', content) if s.strip()][:3]
                return lines[:3]
    except Exception as e:
        logging.warning(f"Reja olish xatosi: {e}")
        return [topic, topic, topic]

async def ask_ai_subtopics(topic: str, main_bands: list, needed: int, lan: str) -> list:
    """Kerakli miqdorda subtopic (raqamsiz)."""
    if needed <= 0:
        return []
    mb = "; ".join(main_bands)
    if lan == "uz":
        prompt = (f"'{topic}' mavzusi uchun {needed} ta qisqa (1-6 so'z) alohida subtopic yozing, "
                  f"ular {mb} (asosiy 3ta) ga takrorlanmasin. Har bir yangi qatorda bitta subtopic bo'lsin. Raqamsiz.")
    elif lan == "ru":
        prompt = (f"–î–ª—è —Ç–µ–º—ã '{topic}' —Å –æ—Å–Ω–æ–≤–Ω—ã–º–∏ –ø—É–Ω–∫—Ç–∞–º–∏: {mb} —Å–≥–µ–Ω–µ—Ä–∏—Ä—É–π—Ç–µ {needed} –∫—Ä–∞—Ç–∫–∏—Ö –ø–æ–¥—Ç–µ–º (1-6 —Å–ª–æ–≤), "
                  "–∫–æ—Ç–æ—Ä—ã–µ –Ω–µ –ø–æ–≤—Ç–æ—Ä—è—é—Ç –æ—Å–Ω–æ–≤–Ω—ã–µ –ø—É–Ω–∫—Ç—ã. –í–µ—Ä–Ω–∏—Ç–µ –ø–æ –æ–¥–Ω–æ–π —Å—Ç—Ä–æ–∫–µ –Ω–∞ –ø–æ–¥—Ç–µ–º—É, –±–µ–∑ –Ω—É–º–µ—Ä–∞—Ü–∏–∏.")
    else:
        prompt = (f"For topic '{topic}' with main points: {mb}, generate {needed} short subtopics (1-6 words) "
                  "that do not duplicate the main points. Return one per line, no numbering.")

    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    try:
        async with aiohttp.ClientSession() as session:
            data_json = {"model": "qwen-turbo", "messages": [{"role": "user", "content": prompt}], "temperature": 0.7}
            async with session.post(API_URL, headers=headers, json=data_json) as resp:
                result = await resp.json()
                content = result["choices"][0]["message"]["content"]
                lines = [clean_title_line(l) for l in content.split('\n') if l.strip()]
                if len(lines) < needed:
                    # Zarurat bo‚Äòlsa, asosiy bandlardan hosila qo‚Äòshamiz
                    extra = []
                    for m in main_bands:
                        extra.append(f"{m} ‚Äî qo‚Äòshimcha")
                        if len(lines) + len(extra) >= needed:
                            break
                    lines.extend(extra)
                return lines[:needed]
    except Exception as e:
        logging.warning(f"Subtopic olish xatosi: {e}")
        return []

def clean_title_line(line: str) -> str:
    """Raqam va belgilarni yo‚Äòqotish."""
    cleaned = re.sub(r'^\s*\d+[\.\)\-\:]\s*', '', line).strip()
    return cleaned

def format_text(text: str, lan: str, line_length: int = 60) -> str:
    lines = []
    for paragraph in text.split('\n'):
        if paragraph.strip():
            wrapped = textwrap.fill(paragraph, width=line_length)
            lines.append(wrapped)
    return '\n'.join(lines)

# -------- SLIDE HELPERS --------
def add_title_slide(bg_num, prs, bg_bytes, title, author, lan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Agar fon rasmi mavjud bo‚Äòlsa ‚Äî qo‚Äòshamiz
    if bg_bytes:
        try:
            slide.shapes.add_picture(BytesIO(bg_bytes), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        except Exception as e:
            logging.warning(f"Fon rasm qo‚Äòshishda xato, fonsiz ketadi: {e}")
    white_text = bg_num in (6, 7, 8)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8),
                                         prs.slide_width - Inches(2.4), Inches(1.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = 'Amasis MT Pro Black'
    p.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)

    # Author
    author_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.4),
                                          prs.slide_width - Inches(2.4), Inches(1.0))
    af = author_box.text_frame
    af.clear()
    ap = af.paragraphs[0]
    if lan == 'uz':
        ap.text = f"Muallif: {author}"
    elif lan == 'ru':
        ap.text = f"–ê–≤—Ç–æ—Ä: {author}"
    else:
        ap.text = f"Author: {author}"
    ap.alignment = PP_ALIGN.CENTER
    ap.font.size = Pt(28)
    ap.font.name = 'Amasis MT Pro Black'
    ap.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)

def add_slide_with_side_image(bg_num, prs, bg_bytes, title, lan, content="", image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_bytes:
        try:
            slide.shapes.add_picture(BytesIO(bg_bytes), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        except Exception as e:
            logging.warning(f"Fon rasm qo‚Äòshishda xato (slide): {e}")
    white_text = bg_num in (6, 7, 8)

    if image_path and os.path.exists(image_path):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5),
                                             prs.slide_width/2 - Inches(1.6), Inches(1))
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7),
                                               prs.slide_width/2 - Inches(1.6),
                                               prs.slide_height - Inches(1.8))
        content_align = PP_ALIGN.LEFT
    else:
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0),
                                             prs.slide_width - Inches(1.6), Inches(1))
        content_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.5),
                                               prs.slide_width - Inches(4.0),
                                               prs.slide_height - Inches(4.0))
        content_align = PP_ALIGN.CENTER

    # Title
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = 'Times New Roman'
    p.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER if not image_path else PP_ALIGN.LEFT

    # Content
    tfc = content_box.text_frame
    tfc.clear()
    if content:
        formatted = format_text(content, lan, line_length=55 if image_path else 80)
        first = True
        for line in formatted.split('\n'):
            if not line.strip():
                continue
            para = tfc.paragraphs[0] if first else tfc.add_paragraph()
            para.text = line
            para.font.size = Pt(16)
            para.font.name = 'Times New Roman'
            para.font.color.rgb = RGBColor(255, 255, 255) if white_text else RGBColor(0, 0, 0)
            para.alignment = content_align
            first = False

    # Rasm
    if image_path and os.path.exists(image_path):
        pic_width = prs.slide_width / 2 - Inches(1.0)
        pic_left = prs.slide_width - pic_width - Inches(0.8)
        pic_top = Inches(1.7)
        try:
            slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width)
        except Exception as e:
            logging.warning(f"Rasm qo‚Äòshishda xato: {e}")


# ==============================
# ===== PRESENTATION CORE ======
# ==============================
async def generate_presentation(uid: int, message: Message, lan):
    """
    - 3 ta asosiy band + kerakli bo‚Äòlsa subtopiclar
    - Rasm rejimi:
        user  -> faqat user yuborgan rasmlar (hech qachon auto qidirmaydi)
        auto  -> Google‚Äôdan avtomatik
        none  -> rasmsiz
    - Har 3-slaydda rasm (0,3,6,‚Ä¶).
    """
    data = user_data.get(uid, {})
    topic = data.get("topic", "Mavzu")
    author = data.get("author", "")

    bg_num = user_data[uid].get("bg_num")
    print(bg_num,'bg_numa')

    bg = data.get("bg", None)  # bytes
    image_mode = data.get("image_mode")  # 'user' | 'auto' | 'none'
    print(image_mode,'image mode')
    requested_slide_count = data.get("requested_slide_count", 20)
    db.add_shablon_pre(
        tg_id=uid,
        bg_num=bg_num,
        ism_fam=author,
        slayid_soni=requested_slide_count,
        til='uz',
        pre_tili=lan
    )
    try:
        requested_slide_count = int(requested_slide_count)
        if requested_slide_count <= 0:
            requested_slide_count = 20
    except:
        requested_slide_count = 20

    # 1) main 3 band
    main_bands = await ask_ai_generate_three(topic, lan)

    # 2) kerakli bo‚Äòlsa subtopics
    remaining = requested_slide_count - len(main_bands)
    subtopics = []
    if remaining > 0:
        subtopics = await ask_ai_subtopics(topic, main_bands, remaining, lan)

    # Slayd sarlavhalari
    full_titles = []
    for i in range(requested_slide_count):
        if i < len(main_bands):
            full_titles.append(main_bands[i])
        else:
            idx_sub = i - len(main_bands)
            if idx_sub < len(subtopics):
                full_titles.append(subtopics[idx_sub])
            else:
                full_titles.append(main_bands[i % len(main_bands)])

    # Save reja
    user_data.setdefault(uid, {})["reja"] = main_bands
    user_data[uid]["requested_slide_count"] = requested_slide_count

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    add_title_slide(bg_num, prs, bg, topic, author, lan)

    # Reja slide
    reja_text = "\n".join([f"{i+1}. {b}" for i, b in enumerate(main_bands)])
    add_slide_with_side_image(bg_num, prs, bg, ("–ü–ª–∞–Ω" if lan=='ru' else ("Reja" if lan=='uz' else "Plan")), lan, reja_text, image_path=None)

    # Progress
    progress_msg = await message.answer("‚è≥ Tayyorlanmoqda...")

    # User rasmlar indexi
    if image_mode == "user":
        user_data[uid]["user_img_idx"] = 0
        local_user_imgs = list(user_images.get(uid, []))
    else:
        local_user_imgs = []

    images_cache = {}  # auto rejimda band bo‚Äòyicha cache

    # Slaydlar
    for idx in range(requested_slide_count):
        try:
            slide_title = f"{idx+1}. {clean_title_line(full_titles[idx])}"

            # Qaysi slaydlarga rasm beramiz?
            wants_image = (idx % 3 == 0) and (image_mode != "none")

            image_path = None
            if wants_image:
                if image_mode == "user":
                    # Faqat foydalanuvchi yuborgan rasmlar
                    if local_user_imgs:
                        # Navbat bilan qo‚Äòyamiz
                        img_idx = user_data[uid].get("user_img_idx", 0)
                        if img_idx < len(local_user_imgs):
                            image_path = local_user_imgs[img_idx]
                            user_data[uid]["user_img_idx"] = img_idx + 1
                        else:
                            image_path = None  # rasmlar tugadi -> boshqa rasm qo‚Äòshmaymiz
                    else:
                        image_path = None
                elif image_mode == "auto":
                    band_for_image = f"{topic} {full_titles[idx]}"
                    if band_for_image in images_cache:
                        image_path = images_cache[band_for_image]
                    else:
                        imgs = download_images(band_for_image, num_images=1)
                        image_path = imgs[0] if imgs else None
                        images_cache[band_for_image] = image_path
                else:
                    image_path = None  # none rejim

            # Matn
            try:
                matn = await ask_ai_content(topic, full_titles[idx], lan=lan)
            except Exception:
                matn = await ask_ai_content_simple(topic, full_titles[idx], lan)

            if image_path:
                # qisqartiramiz (chap kolonka)
                short = matn[:420]
                last = max(short.rfind('.'), short.rfind('!'), short.rfind('?'))
                if last > 0:
                    short = short[:last+1]
                matn = short

            add_slide_with_side_image(bg_num, prs, bg, slide_title, lan, matn, image_path=image_path)

            # progress
            try:
                progress = (idx + 1) * 100 // requested_slide_count
                blocks = max(1, min(10, (progress // 10)))
                bar = ''.join(['üü©' if i < blocks else '‚¨úÔ∏è' for i in range(10)])
                await progress_msg.edit_text(f"‚è≥ Tayyorlanmoqda...\n\n{bar} ({progress}%)")
            except Exception:
                pass

        except Exception as e:
            logging.error(f"Xatolik slayd {idx+1}: {e}")
            try:
                fallback_text = await ask_ai_content_simple(topic, full_titles[idx], lan)
            except Exception:
                fallback_text = "Matn mavjud emas."
            add_slide_with_side_image(bg_num, prs, bg, slide_title, lan, fallback_text, image_path=None)
            continue

    # Save and send
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    caption = f"üìå {topic}\nüí° @Slaydai_bot"
    await message.answer_document(
        BufferedInputFile(buffer.read(), filename=f"{topic}.pptx"),
        caption=caption
    )



# ==========================
# ======  HANDLERS  ========
# ==========================

@router.message(F.text=='/new')
@router.callback_query(F.data=='taqdimot')
async def start(callback: CallbackQuery, state: FSMContext):
    olish=db.select_shablon_pre(tg_id=callback.from_user.id,til='uz')
    if olish:

        pul = db.select_summa()
        summa = db.select_user(tg_user=callback.from_user.id,til='uz')
        tekshir = int(pul[1]) * 5
        if int(summa[3]) >= tekshir:
            try:
                await callback.message.delete()
            except:
                pass

            await callback.message.answer("Taqdimot mavzusini kiriting:")
            await state.set_state(PresState_shablon.topic)
        else:
            text = (f"üí∞Balansingiz: {pul[1]} so'm\n\n"
                    f"Xisobingizni to'ldirish uchun /bye buyrug'idan foydalaning.")
            await callback.message.answer(text, parse_mode="Markdown")
    else:

        pul = db.select_summa()
        summa = db.select_user(tg_user=callback.from_user.id)
        tekshir = int(pul[1]) * 5
        if int(summa[3]) >= tekshir:
            try:
                await callback.message.delete()
            except:
                pass

            keyboard = InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="üá∫üáø O'zbekcha", callback_data="lang_uz")],
                [InlineKeyboardButton(text="üá∑üá∫ –†—É—Å—Å–∫–∏–π", callback_data="lang_ru")],
                [InlineKeyboardButton(text="üá¨üáß English", callback_data="lang_en")]
            ])
            await callback.message.answer("Taqdimot tilini tanlang:", reply_markup=keyboard)
            await state.set_state(PresState.language)
        else:
            text = (f"üí∞Balansingiz: {pul[1]} so'm\n\n"
                    f"Xisobingizni to'ldirish uchun /bye buyrug'idan foydalaning.")
            await callback.message.answer(text, parse_mode="Markdown")

@router.callback_query(StateFilter(PresState.language), F.data.startswith("lang_"))
async def set_language(callback: CallbackQuery, state: FSMContext):
    lang = callback.data.split("_")[1]
    await state.update_data({'lang': lang})
    try:
        await callback.message.delete()
    except:
        pass
    await callback.message.answer("üìå Taqdimot mavzusini kiriting:")
    await state.set_state(PresState.topic)

@router.message(StateFilter(PresState.topic))
async def get_topic(msg: Message, state: FSMContext):
    user_data[msg.from_user.id] = {"topic": msg.text}
    await msg.answer("‚úçÔ∏è Muallif ismini kiriting:")
    await state.set_state(PresState.author)

@router.message(StateFilter(PresState.author))
async def get_author(msg: Message, state: FSMContext):
    user_data[msg.from_user.id]["author"] = msg.text

    # Background selection buttons
    inline_buttons = InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="1", callback_data="bg_1"),
            InlineKeyboardButton(text="2", callback_data="bg_2"),
            InlineKeyboardButton(text="3", callback_data="bg_3"),
            InlineKeyboardButton(text="4", callback_data="bg_4"),
            InlineKeyboardButton(text="5", callback_data="bg_5")
        ],
        [
            InlineKeyboardButton(text="6", callback_data="bg_6"),
            InlineKeyboardButton(text="7", callback_data="bg_7"),
            InlineKeyboardButton(text="8", callback_data="bg_8"),
            InlineKeyboardButton(text="9", callback_data="bg_9"),
            InlineKeyboardButton(text="10", callback_data="bg_10")
        ],
    ])

    rasim_manzili = Path(r".\handlers\users\tanlov.jpg")
    photo = FSInputFile(rasim_manzili)
    await bot.send_photo(
        chat_id=msg.from_user.id,
        photo=photo,
        caption="üì∏ Taqdimot fonini tanlang (1-10) yoki o‚Äòzingiz rasm yuboring.",
        reply_markup=inline_buttons
    )
    await state.set_state(PresState.bg)

@router.callback_query(StateFilter(PresState.bg), F.data.startswith("bg_"))
async def handle_bg_selection(callback: CallbackQuery, state: FSMContext):
    bg_number = callback.data.split("_")[1]
    await state.update_data({"bg_num": bg_number})
    user_data.setdefault(callback.from_user.id, {})["bg_num"] = int(bg_number)
    try:
        rasim_manzili = Path(rf".\handlers\users\rasimlar\{bg_number}.jpg")
        with open(rasim_manzili, "rb") as f:
            user_data[callback.from_user.id]["bg"] = f.read()

        try:
            await callback.message.delete()
        except:
            pass
        await state.set_state(PresState.reja_mode)
        await callback.message.answer("üìÑ Nechta slayd (faqat raqam):")
        await state.set_state(PresState.slide_count)
        await state.update_data(reja_mode="ai")
    except Exception as e:
        logging.error(f"Background load error: {e}")
        await callback.answer("‚ùå Fon yuklanmadi, qayta urining", show_alert=True)

@router.message(StateFilter(PresState.bg_upload), F.photo)
async def handle_uploaded_bg(msg: Message, state: FSMContext):
    photo = msg.photo[-1]
    file = await bot.get_file(photo.file_id)
    bg = await bot.download_file(file.file_path)
    user_data.setdefault(msg.from_user.id, {})["bg"] = bg.read()

    await state.set_state(PresState.reja_mode)
    await msg.answer("üìÑ Nechta slayd (faqat raqam):")
    await state.set_state(PresState.slide_count)
    await state.update_data(reja_mode="ai")

@router.message(StateFilter(PresState.slide_count))
async def get_slide_count(msg: Message, state: FSMContext):
    summa = db.select_summa()
    try:
        if not msg.text.isdigit():
            await msg.answer("‚ùå Iltimos, faqat raqam kiriting!")
            return

        requested_slide_count = int(msg.text)
        if requested_slide_count < 5:
            await msg.answer("‚ùå Slaydlar soni eng kami 5 ta bo'lishi kerak!")
            return
        if requested_slide_count > 50:
            await msg.answer("‚ùå Slaydlar soni 50 tadan ko'p bo'lishi mumkin emas!")
            return

        check = db.select_user(tg_user=msg.from_user.id)
        tek = requested_slide_count * int(summa[1])
        if int(check[3]) < tek:
            await msg.answer(
                f"Xisobingizda yetarli mablag' yo'q üôÖ!\n"
                f"Xisobingiz:üí∞ {check[3]}\n\n"
                f"Xisobni to‚Äòldirish uchun /bye buyrug‚Äòidan foydalaning."
            )
            await state.clear()
            return

        # Pul yechish
        update = int(check[3]) - tek
        db.update_user_balans(tg_user=msg.from_user.id, balans=update)

        uid = msg.from_user.id
        data = await state.get_data()
        lan = data.get('lang', 'uz')
        user_data.setdefault(uid, {})["requested_slide_count"] = requested_slide_count

        # Tugmalar
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="‚úÖ Rasm tanlash", callback_data="img_select")],
            [InlineKeyboardButton(text="‚ôªÔ∏è Avtomatik rasmlar", callback_data="img_auto")],
            [InlineKeyboardButton(text="üö´ Rasmsiz", callback_data="img_none")]
        ])

        a=await msg.answer(
            "üìÑ Slaydlar soni qabul qilindi.\n\n"
            "‚úÖ *Rasm tanlash* ‚Äî mavzuga oid rasmlarni o‚Äòzingiz yuborasiz (7 tagacha) va faqat shu rasmlar ishlatiladi.\n"
            "‚ôªÔ∏è *Avtomatik rasmlar* ‚Äî bot mavzuga oid rasmlarni internetdan yuklaydi.\n"
            "üö´ *Rasmsiz* ‚Äî taqdimot rasm ishlatmaydi.",
            reply_markup=kb,
            parse_mode="Markdown"
        )
        await state.update_data({"m_id":a.message_id})
        await state.update_data(reja_mode="ai")

    except Exception as e:
        logging.error(f"get_slide_count xato: {e}")
        await state.clear()

# --- Rasm tanlash rejimi ---
@router.callback_query(F.data == "img_select")
async def select_images(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol.get("m_id")
    await bot.delete_message(chat_id=callback.message.chat.id, message_id=m_id)
    uid = callback.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "user"
    clear_user_images(uid, delete_files=True)

    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="‚úÖ Tugatish", callback_data="img_finish")],
        [InlineKeyboardButton(text="‚ôª Rasmlarni o'chirish", callback_data="img_reset")]
    ])
    await callback.message.answer(
        "üì§ Mavzuga oid suratlarni yuboring (7 tagacha).\n"
"Tayyor bo'lgach, *Finish* tugmasini bosing.\n"
"Agar xato qilsangiz ‚Äî *Clear Photos* orqali qaytadan boshlang.\n"
"üåÉ Rasmlar: 0/7",
        reply_markup=kb,
        parse_mode="Markdown"
    )
    await state.set_state(ImageCollectState.collecting)

@router.message(StateFilter(ImageCollectState.collecting), F.photo)
async def collect_images(msg: Message, state: FSMContext):
    uid = msg.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "user"
    if uid not in user_images:
        user_images[uid] = []
    if len(user_images[uid]) >= 7:
        await msg.answer("‚ùå Siz allaqachon 7 ta rasm yuborgansiz!")
        return

    file = await bot.get_file(msg.photo[-1].file_id)
    img = await bot.download_file(file.file_path)
    tmp_path = tempfile.mktemp(suffix=".jpg")
    with open(tmp_path, "wb") as f:
        f.write(img.read())

    user_images[uid].append(tmp_path)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="‚úÖ Tugatish", callback_data="img_finish")],
        [InlineKeyboardButton(text="‚ôª Udalit fotografii", callback_data="img_reset")]
    ])
    await msg.answer(f"üåÉ Rasmlar: {len(user_images[uid])}/7", reply_markup=kb)

@router.callback_query(F.data == "img_reset")
async def reset_images(callback: CallbackQuery):

    uid = callback.from_user.id
    # barcha yuborgan rasmlarini bekor qilamiz (fayllarni ham o‚Äòchiramiz)
    clear_user_images(uid, delete_files=True)
    await callback.message.answer(
        "‚ôª 0 ta surat yuklandi. Qayta yuborishingiz mumkin.\nüåÉ Rasmlar: 0/7"
    )

@router.callback_query(F.data == "img_finish")
async def finish_images(callback: CallbackQuery, state: FSMContext):
    uid = callback.from_user.id
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="‚úÖ Tayyorgarlik", callback_data="tayyorlash"),
                InlineKeyboardButton(text="üö´ Rad etish", callback_data="rad_etish"),
                InlineKeyboardButton(text="‚úèÔ∏è O'zgartirish", callback_data="ozgartirish")
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a = await bot.send_photo(chat_id=callback.from_user.id, caption=
    "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"Mavzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
f"‚Ä¢ Agar ma'lumot to'g'ri bo'lsa, '‚úÖ Tayyorlash' tugmasini bosing!\n"
f"‚Ä¢ O'zgartirish uchun '‚úèÔ∏è Tahrirlash' tugmasini bosing\n"
f"‚Ä¢ Bekor qilish uchun 'üö´ Rad etish' tugmasini bosing",
                             reply_markup=kb, photo=rasim
                             )
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState.tanlov)
    await callback.answer()

# --- Rasmsiz rejim ---
@router.callback_query(F.data == "img_none")
async def no_images(callback: CallbackQuery, state: FSMContext):

    uid = callback.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "none"
    clear_user_images(uid, delete_files=True)  # albatta ishlatilmaydi
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="‚úÖ Tayyorgarlik", callback_data="tayyorlash"),
                InlineKeyboardButton(text="üö´ Rad etish", callback_data="rad_etish"),
                InlineKeyboardButton(text="‚úèÔ∏è O'zgartirish", callback_data="ozgartirish")
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a = await bot.send_photo(chat_id=callback.from_user.id, caption=
    "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"Mavzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
f"‚Ä¢ Agar ma'lumot to'g'ri bo'lsa, '‚úÖ Tayyorlash' tugmasini bosing!\n"
f"‚Ä¢ O'zgartirish uchun '‚úèÔ∏è Tahrirlash' tugmasini bosing\n"
f"‚Ä¢ Bekor qilish uchun 'üö´ Rad etish' tugmasini bosing",
                             reply_markup=kb, photo=rasim
                             )
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState.tanlov)
    await callback.answer()





# --- Avtomatik rasm rejimi ---
@router.callback_query(F.data == "img_auto")
async def img_auto_handler(callback: CallbackQuery, state: FSMContext):

    uid = callback.from_user.id
    user_data.setdefault(uid, {})["image_mode"] = "auto"
    clear_user_images(uid, delete_files=True)

    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")


    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="‚úÖ Tayyorgarlik", callback_data="tayyorlash"),
                InlineKeyboardButton(text="üö´ Rad etish", callback_data="rad_etish"),
                InlineKeyboardButton(text="‚úèÔ∏è O'zgartirish", callback_data="ozgartirish")
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a=await bot.send_photo(chat_id=callback.from_user.id,caption=
        "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"Mavzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
f"‚Ä¢ Agar ma'lumot to'g'ri bo'lsa, '‚úÖ Tayyorlash' tugmasini bosing!\n"
f"‚Ä¢ O'zgartirish uchun '‚úèÔ∏è Tahrirlash' tugmasini bosing\n"
f"‚Ä¢ Bekor qilish uchun 'üö´ Rad etish' tugmasini bosing",
        reply_markup=kb,photo=rasim
    )
    await state.update_data({"m_id":a.message_id})
    await state.set_state(PresState.tanlov)
    await callback.answer()

@router.callback_query(F.data=="mavzu",PresState.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    a = await callback.message.answer("Taqdimot mavzusini yuboring:")
    await state.update_data({"mm_id": a.message_id})

    await state.set_state(PresState.mavzu)
    await callback.answer()
    print(2)
@router.callback_query(F.data=="til",PresState.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="üá∫üáø O'zbekcha", callback_data="uz")],
        [InlineKeyboardButton(text="üá∑üá∫ –†—É—Å—Å–∫–∏–π", callback_data="ru")],
        [InlineKeyboardButton(text="üá¨üáß English", callback_data="en")]
    ])
    a = await callback.message.answer("Taqdimot tilini tanlang:", reply_markup=keyboard)
    await state.update_data({"mm_id": a.message_id})
    await state.set_state(PresState.til)
    await callback.answer()
    print(1)


@router.callback_query(F.data=="muallif",PresState.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):

        a = await callback.message.answer("Muallif nomini yuboring:")
        await state.update_data({"mm_id": a.message_id})
        await state.set_state(PresState.muallif)
        await callback.answer()
        print(3)

@router.callback_query(F.data == "ozgartirish", PresState.tanlov)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol.get("m_id")
    uid=callback.from_user.id
    await bot.delete_message(chat_id=callback.from_user.id, message_id=m_id)
    data = await state.get_data()
    lan = data.get('lang', 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")
    lang = data.get('lang')
    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")


    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu"),
                InlineKeyboardButton(text="Muallif ", callback_data="muallif"),
                InlineKeyboardButton(text="Til", callback_data="til")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus"),
                InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}", callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus")
            ],
            [
                InlineKeyboardButton(text="‚¨ÖÔ∏è", callback_data="oldingi"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="‚û°Ô∏è", callback_data="keyingi")
            ],
            [
                InlineKeyboardButton(text="Save ‚úÖ", callback_data="saqlash"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    rasim = FSInputFile(photo_manzili)
    a=await bot.send_photo(chat_id=callback.from_user.id,caption=
        "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"Mavzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
,
        reply_markup=kb,photo=rasim
    )
    await state.update_data({"m_id":a.message_id})
    await state.set_state(PresState.ozgartirish)




@router.callback_query(PresState.ozgartirish)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):

    data=callback.data

    if data=='saqlash':
        uid = callback.from_user.id

        ol = await state.get_data()
        m_id = ol.get("m_id")
        await bot.delete_message(chat_id=callback.from_user.id, message_id=m_id)

        data = await state.get_data()
        lan = data.get('lang', 'uz')

        author = user_data[uid].get("author")

        requested_slide_count = user_data[uid].get("requested_slide_count", 0)
        bg_num = user_data[uid].get("bg_num")
        await generate_presentation(uid, callback.message, lan)



        await state.clear()

    else:
        ol=await state.get_data()
        m_id=ol.get("m_id")
        bg_num = user_data[callback.from_user.id].get("bg_num")
        requested_slide_count = user_data[callback.from_user.id].get("requested_slide_count", 0)
        if data == "plus":
            if int(requested_slide_count) > 19:
                await callback.answer("Siz maksimal 20 sahifani tanlashingiz mumkin.‚ùóÔ∏è", show_alert=True)
            else:
                user_data.setdefault(callback.from_user.id, {})["requested_slide_count"] = int(
                    requested_slide_count) + 1
        if data == "minus":
            if int(requested_slide_count) < 6:
                await callback.answer("Siz maksimal 5 sahifani tanlashingiz mumkin.‚ùóÔ∏è", show_alert=True)
            else:
                user_data.setdefault(callback.from_user.id, {})["requested_slide_count"] = int(
                    requested_slide_count) - 1
        if data == "oldingi":
            if int(bg_num) == 1:
                user_data.setdefault(callback.from_user.id, {})["bg_num"] = 10
                rasim_manzili = Path(rf".\handlers\users\rasimlar\{10}.jpg")
                with open(rasim_manzili, "rb") as f:
                    user_data[callback.from_user.id]["bg"] = f.read()
            else:
                user_data.setdefault(callback.from_user.id, {})["bg_num"] = int(bg_num) - 1
                rasim_manzili = Path(rf".\handlers\users\rasimlar\{bg_num - 1}.jpg")
                with open(rasim_manzili, "rb") as f:
                    user_data[callback.from_user.id]["bg"] = f.read()
        if data == "keyingi":
            if int(bg_num) == 10:

                user_data.setdefault(callback.from_user.id, {})["bg_num"] = 1
                rasim_manzili = Path(rf".\handlers\users\rasimlar\{1}.jpg")
                with open(rasim_manzili, "rb") as f:
                    user_data[callback.from_user.id]["bg"] = f.read()
            else:
                user_data.setdefault(callback.from_user.id, {})["bg_num"] = int(bg_num) + 1
                rasim_manzili = Path(rf".\handlers\users\rasimlar\{bg_num + 1}.jpg")
                with open(rasim_manzili, "rb") as f:
                    user_data[callback.from_user.id]["bg"] = f.read()


        ol = await state.get_data()
        m_id = ol.get("m_id")
        uid = callback.from_user.id

        data = await state.get_data()
        lan = data.get('lang', 'uz')
        topic = user_data[uid].get("topic")
        author = user_data[uid].get("author")
        lang = data.get('lang')
        requested_slide_count = user_data[uid].get("requested_slide_count", 0)
        bg_num = user_data[uid].get("bg_num")

        kb = InlineKeyboardMarkup(
            inline_keyboard=[
                [
                    InlineKeyboardButton(text="Mavzu", callback_data="mavzu"),
                    InlineKeyboardButton(text="Muallif", callback_data="muallif"),
                        InlineKeyboardButton(text="Til", callback_data="til")
                ],
                [
                    InlineKeyboardButton(text="-", callback_data="minus"),
                    InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}",
                                         callback_data=f"count:{requested_slide_count}"),
                    InlineKeyboardButton(text="+", callback_data="plus")
                ],
                [
                    InlineKeyboardButton(text="‚¨ÖÔ∏è", callback_data="oldingi"),
                    InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                    InlineKeyboardButton(text="‚û°Ô∏è", callback_data="keyingi")
                ],
                [
                    InlineKeyboardButton(text="Saqlash ‚úÖ", callback_data="saqlash"),
                ]
            ]
        )
        photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")

        rasim = FSInputFile(photo_manzili)
        media = InputMediaPhoto(
            media=rasim,caption=
        "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"Mavzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
)

        a = await bot.edit_message_media(chat_id=callback.from_user.id,media=media,message_id=m_id,reply_markup=kb)
        await state.update_data({"m_id": a.message_id})
        await state.set_state(PresState.ozgartirish)






@router.callback_query(F.data == "rad_etish", PresState.tanlov)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    await bot.send_message(chat_id=callback.from_user.id,text="Taqdimot bekor qilindi.\n"
"Yangi taqdimot yaratish uchun quyidagi üìïTaqdimot tugmasini bosing!\n\n"
"Davom etish uchun /start tugmasini bosing.")

    await state.clear()
    await callback.answer()

@router.callback_query(F.data == "tayyorlash", PresState.tanlov)
async def tayyorlash_handler(callback: CallbackQuery, state: FSMContext):
    ol=await state.get_data()
    m_id=ol["m_id"]
    await bot.delete_message(chat_id=callback.from_user.id,message_id=m_id)
    uid = callback.from_user.id


    data = await state.get_data()
    lan = data.get('lang', 'uz')

    author = user_data[uid].get("author")

    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")
    await generate_presentation(uid, callback.message, lan)



    await state.clear()


@router.callback_query(StateFilter(PresState.til))
async def handle_mavzu_message(message: CallbackQuery, state: FSMContext):
    """
    Handle language selection for presentation template (English version).
    Updates the language in database and refreshes the interface.
    """
    # Extract data and user info
    selected_lang = message.data
    user_id = message.from_user.id

    # Get previous message ID and delete it
    state_data = await state.get_data()
    old_msg_id = state_data.get("mm_id")
    await bot.delete_message(chat_id=user_id, message_id=old_msg_id)

    # Update language in database
    db.update_shablon_lang(
        tg_id=user_id,
        pre_tili='en',  # Previous language was English
        til=selected_lang
    )

    # Update user data
    user_data.setdefault(user_id, {})["lang"] = selected_lang

    # Prepare updated interface data
    state_data = await state.get_data()
    main_msg_id = state_data.get("m_id")
    topic = user_data[user_id].get("topic")
    author = user_data[user_id].get("author")
    requested_slide_count = user_data[user_id].get("requested_slide_count", 0)
    bg_num = user_data[user_id].get("bg_num")

    # Create keyboard
    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu"),
                InlineKeyboardButton(text="Muallif", callback_data="muallif"),
                InlineKeyboardButton(text="Til", callback_data="til")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus"),
                InlineKeyboardButton(
                    text=f"Sahifalar: {requested_slide_count}",
                    callback_data=f"count:{requested_slide_count}"
                ),
                InlineKeyboardButton(text="+", callback_data="plus")
            ],
            [
                InlineKeyboardButton(text="‚¨ÖÔ∏è", callback_data="oldingi"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="‚û°Ô∏è", callback_data="keyingi")
            ],
            [
                    InlineKeyboardButton(text="Saqlash ‚úÖ", callback_data="saqlash"),
            ]
        ]
    )

    # Prepare media message
    photo_path = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")
    image = FSInputFile(photo_path)

    media = InputMediaPhoto(
        media=image,
        caption=(
            "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
            f"Mavzu: {topic}\n"
            f"Muallif: {author}\n"
            f"Sahifalar soni: {requested_slide_count}\n"
            f"Til: {selected_lang}\n\n"

        )
    )

    # Update the message
    updated_msg = await bot.edit_message_media(
        chat_id=user_id,
        media=media,
        message_id=main_msg_id,
        reply_markup=kb
    )

    # Update state
    await state.update_data({"m_id": updated_msg.message_id})
    await state.set_state(PresState.ozgartirish)



@router.message(StateFilter(PresState.muallif))
async def handle_mavzu_message(message: Message, state: FSMContext):
    print(1,'011')
    ol = await state.get_data()
    mm_id=ol.get("mm_id")
    await bot.delete_message(chat_id=message.from_user.id,message_id=mm_id)
    await message.delete()
    user_data.setdefault(message.from_user.id, {})["author"] = message.text
    ol = await state.get_data()
    m_id = ol.get("m_id")
    uid = message.from_user.id

    data = await state.get_data()
    lan = user_data[uid].get("lang", 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")

    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu"),
                InlineKeyboardButton(text="Muallif ", callback_data="muallif"),
                InlineKeyboardButton(text="Til", callback_data="til")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus"),
                InlineKeyboardButton(text=f"Sahifalar: {requested_slide_count}",
                                     callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus")
            ],
            [
                InlineKeyboardButton(text="‚¨ÖÔ∏è", callback_data="oldingi"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="‚û°Ô∏è", callback_data="keyingi")
            ],
            [
                InlineKeyboardButton(text="Saqlash ‚úÖ", callback_data="saqlash"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")

    rasim = FSInputFile(photo_manzili)
    media = InputMediaPhoto(
        media=rasim, caption=
        "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"Mavzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
)

    a = await bot.edit_message_media(chat_id=message.from_user.id, media=media, message_id=m_id, reply_markup=kb)
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState.ozgartirish)





@router.message(StateFilter(PresState.mavzu))
async def handle_mavzu_message(message: Message, state: FSMContext):

    ol = await state.get_data()
    mm_id=ol.get("mm_id")
    await bot.delete_message(chat_id=message.from_user.id,message_id=mm_id)
    await message.delete()
    user_data.setdefault(message.from_user.id, {})["topic"] = message.text
    ol = await state.get_data()
    m_id = ol.get("m_id")
    uid = message.from_user.id

    data = await state.get_data()
    lan  = user_data[uid].get("lang", 'uz')
    topic = user_data[uid].get("topic")
    author = user_data[uid].get("author")

    requested_slide_count = user_data[uid].get("requested_slide_count", 0)
    bg_num = user_data[uid].get("bg_num")

    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [
                InlineKeyboardButton(text="Mavzu", callback_data="mavzu"),
                InlineKeyboardButton(text="Muallif", callback_data="muallif"),
                InlineKeyboardButton(text="Til", callback_data="til")
            ],
            [
                InlineKeyboardButton(text="-", callback_data="minus"),
                InlineKeyboardButton(text=f"Sahifa: {requested_slide_count}",
                                     callback_data=f"count:{requested_slide_count}"),
                InlineKeyboardButton(text="+", callback_data="plus")
            ],
            [
                InlineKeyboardButton(text="‚¨ÖÔ∏è", callback_data="oldingi"),
                InlineKeyboardButton(text=f"Dizayn {bg_num}", callback_data=f"dizayn:{bg_num}"),
                InlineKeyboardButton(text="‚û°Ô∏è", callback_data="keyingi")
            ],
            [
                InlineKeyboardButton(text="Saqlash ‚úÖ", callback_data="saqlash"),
            ]
        ]
    )
    photo_manzili = Path(rf".\handlers\users\rasimlar\{bg_num}.jpg")

    rasim = FSInputFile(photo_manzili)
    media = InputMediaPhoto(
        media=rasim, caption=
        "üåüAjoyib, quyidagi ma'lumotlarni tekshiring.\n\n"
f"MAvzu: {topic}\n"
f"Muallif: {author}\n"
f"Sahifalar soni: {requested_slide_count}\n"
f"Til: {lan}\n\n"
)

    a = await bot.edit_message_media(chat_id=message.from_user.id, media=media, message_id=m_id, reply_markup=kb)
    await state.update_data({"m_id": a.message_id})
    await state.set_state(PresState.ozgartirish)



# (ixtiyoriy) Reja bilan ishlash (agar kerak bo‚Äòlsa)
@router.message(StateFilter(PresState.reja))
async def reja_input(msg: Message, state: FSMContext):
    if msg.text.lower() == "tayyor":
        if len(user_data[msg.from_user.id].get("reja", [])) == 0:
            await msg.answer("‚ùå Kamida bitta reja bandini kiriting!")
            return
        await msg.answer("üìÑ Nechta slayd (faqat raqam):")
        await state.set_state(PresState.slide_count)
    else:
        user_data[msg.from_user.id].setdefault("reja", []).append(msg.text)

# End of file
